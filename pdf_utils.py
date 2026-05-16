import os
import tempfile
import zipfile
import sys
import subprocess
from PyPDF2 import PdfReader, PdfWriter

try:
    from PIL import Image
except ImportError:
    # Auto-install Pillow in the exact environment running this code
    subprocess.check_call([sys.executable, "-m", "pip", "install", "Pillow"])
    from PIL import Image

IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp'}
WORD_EXTENSIONS = {'.doc', '.docx'}
PPT_EXTENSIONS = {'.ppt', '.pptx'}

def create_pdf_from_image(path):
    with Image.open(path) as img:
        if img.mode in ('RGBA', 'LA') or (img.mode == 'P' and 'transparency' in img.info):
            img = img.convert('RGB')
        else:
            img = img.convert('RGB')
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
        temp_path = temp_file.name
        temp_file.close()
        img.save(temp_path, 'PDF', resolution=100.0)
        return temp_path

def create_pdf_from_word(path):
    import comtypes.client
    word = comtypes.client.CreateObject('Word.Application')
    word.Visible = False
    doc = word.Documents.Open(os.path.abspath(path))
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
    temp_path = temp_file.name
    temp_file.close()
    try:
        doc.SaveAs(temp_path, FileFormat=17)
    finally:
        doc.Close()
        word.Quit()
    return temp_path

def create_pdf_from_ppt(path):
    import comtypes.client
    powerpoint = comtypes.client.CreateObject('PowerPoint.Application')
    pres = powerpoint.Presentations.Open(os.path.abspath(path), WithWindow=False)
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
    temp_path = temp_file.name
    temp_file.close()
    try:
        pres.SaveAs(temp_path, 32)
    finally:
        pres.Close()
        powerpoint.Quit()
    return temp_path

def normalize_to_pdf(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == '.pdf':
        return path
    if ext in IMAGE_EXTENSIONS:
        return create_pdf_from_image(path)
    if ext in WORD_EXTENSIONS:
        return create_pdf_from_word(path)
    if ext in PPT_EXTENSIONS:
        return create_pdf_from_ppt(path)
    raise ValueError(f'Unsupported file type: {ext}. Only PDF, Word, PPT, and image files are supported.')

def merge_pdfs(paths, output):
    writer = PdfWriter()
    file_objs = []
    temp_files = []
    try:
        for path in paths:
            normalized_path = normalize_to_pdf(path)
            if normalized_path != path:
                temp_files.append(normalized_path)
            f = open(normalized_path, 'rb')
            file_objs.append(f)
            reader = PdfReader(f, strict=False)
            for page in reader.pages:
                writer.add_page(page)
        with open(output, 'wb') as out_f:
            writer.write(out_f)
    finally:
        for f in file_objs:
            f.close()
        for temp in temp_files:
            try:
                os.remove(temp)
            except Exception:
                pass

def split_pdf(path, output_zip, split_page):
    normalized_path = normalize_to_pdf(path)
    try:
        with open(normalized_path, 'rb') as f:
            reader = PdfReader(f, strict=False)
            writer1 = PdfWriter()
            writer2 = PdfWriter()
            
            split_page = max(1, min(split_page, len(reader.pages) - 1))
            if len(reader.pages) <= 1:
                raise ValueError('Document only has 1 page, cannot be split into two')
                
            for i in range(split_page):
                writer1.add_page(reader.pages[i])
                
            for i in range(split_page, len(reader.pages)):
                writer2.add_page(reader.pages[i])
                
            part1_path = output_zip.replace('.zip', '_part1.pdf')
            part2_path = output_zip.replace('.zip', '_part2.pdf')
            
            with open(part1_path, 'wb') as out1:
                writer1.write(out1)
            with open(part2_path, 'wb') as out2:
                writer2.write(out2)
                
            with zipfile.ZipFile(output_zip, 'w') as zf:
                zf.write(part1_path, os.path.basename(part1_path))
                zf.write(part2_path, os.path.basename(part2_path))
                
            try:
                os.remove(part1_path)
                os.remove(part2_path)
            except Exception:
                pass
    finally:
        if normalized_path != path:
            try:
                os.remove(normalized_path)
            except Exception:
                pass

def rotate_pdf(path, output, angle):
    normalized_path = normalize_to_pdf(path)
    try:
        with open(normalized_path, 'rb') as f:
            reader = PdfReader(f, strict=False)
            writer = PdfWriter()
            for page in reader.pages:
                try:
                    page.rotate(angle)
                except Exception:
                    try:
                        page.rotate_clockwise(angle)
                    except Exception:
                        pass
                writer.add_page(page)
            with open(output, 'wb') as out_f:
                writer.write(out_f)
    finally:
        if normalized_path != path:
            try:
                os.remove(normalized_path)
            except Exception:
                pass

def remove_pages(path, output, pages_to_remove):
    normalized_path = normalize_to_pdf(path)
    try:
        with open(normalized_path, 'rb') as f:
            reader = PdfReader(f, strict=False)
            writer = PdfWriter()
            pages_to_remove_set = set(pages_to_remove)
            if not pages_to_remove_set:
                raise ValueError('No pages specified to remove')
            for idx, page in enumerate(reader.pages, start=1):
                if idx not in pages_to_remove_set:
                    writer.add_page(page)
            if not writer.pages:
                raise ValueError('Cannot remove all pages from the document')
            with open(output, 'wb') as out_f:
                writer.write(out_f)
    finally:
        if normalized_path != path:
            try:
                os.remove(normalized_path)
            except Exception:
                pass

def extract_text_from_pdf(path):
    normalized_path = normalize_to_pdf(path)
    try:
        with open(normalized_path, 'rb') as f:
            reader = PdfReader(f, strict=False)
            texts = []
            for page in reader.pages:
                try:
                    texts.append(page.extract_text() or "")
                except Exception:
                    texts.append("")
            return "\n".join(texts)
    finally:
        if normalized_path != path:
            try:
                os.remove(normalized_path)
            except Exception:
                pass

def insert_file_into_pdf(target_path, insert_path, output, position):
    normalized_target = normalize_to_pdf(target_path)
    normalized_insert = normalize_to_pdf(insert_path)
    temp_files = []
    if normalized_target != target_path:
        temp_files.append(normalized_target)
    if normalized_insert != insert_path:
        temp_files.append(normalized_insert)
    
    try:
        with open(normalized_target, 'rb') as f_target, open(normalized_insert, 'rb') as f_insert:
            reader_target = PdfReader(f_target, strict=False)
            reader_insert = PdfReader(f_insert, strict=False)
            writer = PdfWriter()
            
            pos_index = min(max(0, position - 1), len(reader_target.pages))
            
            for i in range(pos_index):
                writer.add_page(reader_target.pages[i])
            
            for page in reader_insert.pages:
                writer.add_page(page)
                
            for i in range(pos_index, len(reader_target.pages)):
                writer.add_page(reader_target.pages[i])
                
            with open(output, 'wb') as out_f:
                writer.write(out_f)
    finally:
        for temp in temp_files:
            try:
                os.remove(temp)
            except Exception:
                pass


def pdf_to_word(pdf_path, output_path):
    from docx import Document
    from docx.shared import Pt
    
    text = extract_text_from_pdf(pdf_path)
    doc = Document()
    doc.add_heading('PDF Content', 0)
    for paragraph_text in text.split('\n'):
        if paragraph_text.strip():
            doc.add_paragraph(paragraph_text)
    doc.save(output_path)


def pdf_to_ppt(pdf_path, output_path):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    
    text = extract_text_from_pdf(pdf_path)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "PDF Content"
    subtitle.text = "Converted from PDF"
    
    bullet_slide_layout = prs.slide_layouts[1]
    lines = [l.strip() for l in text.split('\n') if l.strip()]
    
    for i, line in enumerate(lines):
        if i % 5 == 0:
            current_slide = prs.slides.add_slide(bullet_slide_layout)
            title = current_slide.shapes.title
            title.text = f"Content (Part {i // 5 + 1})"
            body_shape = current_slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
        
        if i % 5 == 0:
            tf.text = line[:100]
        else:
            p = tf.add_paragraph()
            p.text = line[:100]
            p.level = 0
    
    prs.save(output_path)
